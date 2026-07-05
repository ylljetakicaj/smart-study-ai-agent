"""
retriever.py
============
Document ingestion + lightweight retrieval (RAG) for the Smart Study AI Agent.

Design decision:
    A full vector database is overkill for single-student lecture-note sets and
    would add heavy dependencies. Instead we use a transparent, dependency-light
    retriever:

      1. Extract text from PDF / PPTX / TXT.
      2. Chunk it into overlapping passages.
      3. Score chunks against a query with TF-IDF cosine similarity.

    This keeps the project easy to run (pure-Python + tiny libs), while the
    Retriever interface is designed so it can be swapped for a Gemini-embeddings
    + vector-store backend without touching the agent code.
"""

from __future__ import annotations

import math
import re
from collections import Counter
from dataclasses import dataclass


# ---------------------------------------------------------------------------
# Text extraction from multiple file types
# ---------------------------------------------------------------------------
def extract_text(file) -> str:
    """
    Extract raw text from a PDF, PPTX, or TXT file.

    Accepts either a file path string or a Streamlit UploadedFile object.
    Parsers are imported lazily to keep startup fast.
    """
    # Streamlit UploadedFile — read from the in-memory buffer
    if hasattr(file, "name"):
        name = file.name.lower()
        data = file.read()
        file.seek(0)  # reset so callers can re-read if needed
        if name.endswith(".pdf"):
            return _extract_pdf_bytes(data)
        if name.endswith((".pptx", ".ppt")):
            return _extract_pptx_bytes(data)
        # Plain text / markdown
        return data.decode("utf-8", errors="ignore")

    # File path string
    lower = str(file).lower()
    if lower.endswith(".pdf"):
        return _extract_pdf_bytes(open(file, "rb").read())
    if lower.endswith((".pptx", ".ppt")):
        return _extract_pptx_bytes(open(file, "rb").read())
    with open(file, encoding="utf-8", errors="ignore") as fh:
        return fh.read()


def _extract_pdf_bytes(data: bytes) -> str:
    from pypdf import PdfReader
    import io

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_pptx_bytes(data: bytes) -> str:
    from pptx import Presentation
    import io

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------
@dataclass
class Chunk:
    text: str
    index: int


def chunk_text(text: str, size: int = 900, overlap: int = 150) -> list[Chunk]:
    """
    Split text into overlapping character windows.

    Overlap preserves context that would otherwise be cut mid-idea at a
    boundary, which improves retrieval quality for questions that straddle
    two chunks.
    """
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks: list[Chunk] = []
    start, idx = 0, 0
    while start < len(text):
        end = start + size
        chunks.append(Chunk(text=text[start:end], index=idx))
        idx += 1
        start += size - overlap  # step forward, keeping the overlap
    return chunks


# ---------------------------------------------------------------------------
# TF-IDF retriever
# ---------------------------------------------------------------------------
_TOKEN_RE = re.compile(r"[a-z0-9]+")


def _tokenize(text: str) -> list[str]:
    return _TOKEN_RE.findall(text.lower())


class Retriever:
    """
    In-memory TF-IDF retriever over the chunks of the current document(s).

    Public interface:
        add_document(text)      -> ingest and index a document
        query(question, k)      -> return the k most relevant chunks
        full_text()             -> the whole corpus (for summarization)
    """

    def __init__(self) -> None:
        self._chunks: list[Chunk] = []
        self._tf: list[Counter] = []          # term frequency per chunk
        self._df: Counter = Counter()         # document frequency per term
        self._raw_text: str = ""

    # ---- ingestion -------------------------------------------------------
    def add_document(self, text: str) -> int:
        """Index a document's text. Returns the number of chunks added."""
        self._raw_text += ("\n" + text) if self._raw_text else text
        new_chunks = chunk_text(text)
        for ch in new_chunks:
            ch.index = len(self._chunks)  # global index
            self._chunks.append(ch)
            counts = Counter(_tokenize(ch.text))
            self._tf.append(counts)
            for term in counts:               # each term counted once per chunk
                self._df[term] += 1
        return len(new_chunks)

    def reset(self) -> None:
        """Clear the index (e.g., when the student uploads a fresh document)."""
        self.__init__()

    # ---- retrieval -------------------------------------------------------
    def _idf(self, term: str) -> float:
        n = max(1, len(self._chunks))
        # Smoothed IDF avoids division-by-zero and dampens rare-term blowups.
        return math.log((1 + n) / (1 + self._df.get(term, 0))) + 1.0

    def _vector(self, counts: Counter) -> dict[str, float]:
        """Build an L2-normalized TF-IDF vector from term counts."""
        vec = {t: (c * self._idf(t)) for t, c in counts.items()}
        norm = math.sqrt(sum(v * v for v in vec.values())) or 1.0
        return {t: v / norm for t, v in vec.items()}

    def query(self, question: str, k: int = 4) -> list[str]:
        """Return the top-k chunk texts most similar to the question."""
        if not self._chunks:
            return []
        q_vec = self._vector(Counter(_tokenize(question)))
        scored: list[tuple[float, int]] = []
        for i, counts in enumerate(self._tf):
            c_vec = self._vector(counts)
            # Cosine similarity = dot product of two normalized vectors.
            shared = set(q_vec) & set(c_vec)
            score = sum(q_vec[t] * c_vec[t] for t in shared)
            scored.append((score, i))
        scored.sort(reverse=True)
        return [self._chunks[i].text for score, i in scored[:k] if score > 0]

    # ---- helpers ---------------------------------------------------------
    def context_for(self, question: str, k: int = 4) -> str:
        """Concatenate the top-k chunks into a single context block."""
        hits = self.query(question, k)
        return "\n\n---\n\n".join(hits) if hits else "(no relevant context found)"

    def full_text(self, max_chars: int = 12000) -> str:
        """Whole corpus, truncated to stay within model limits (summaries)."""
        return self._raw_text[:max_chars]

    @property
    def is_empty(self) -> bool:
        return not self._chunks
